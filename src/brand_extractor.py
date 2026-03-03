"""
Brand Extraction Engine

Extracts structured brand profiles from multiple input types:
- Brand Guide PDFs
- Existing presentations (.pptx)
- Screenshots / images
- Logo files
- Website URLs
- Manual input
"""

import io
import os
import re
import json
import base64
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

import anthropic

from .brand_profile import (
    BrandProfile, ColorPalette, ColorEntry, Typography, FontSpec,
    LayoutSpacing, VisualElements, ContentStyle, Confidence,
    AspectRatio, WhiteSpacePhilosophy, IconStyle, ShapeLanguage,
    WritingTone, HeadlineStyle,
)


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return (r, g, b)


def _rgb_to_hex(r: int, g: int, b: int) -> str:
    return f"#{r:02X}{g:02X}{b:02X}"


def _make_color_entry(hex_val: str, confidence: Confidence = Confidence.MEDIUM,
                      usage: str = "") -> ColorEntry:
    hex_val = hex_val.strip()
    if not hex_val.startswith("#"):
        hex_val = "#" + hex_val
    try:
        rgb = _hex_to_rgb(hex_val)
    except (ValueError, IndexError):
        rgb = (0, 0, 0)
    return ColorEntry(hex=hex_val.upper(), rgb=rgb, confidence=confidence, usage_rule=usage)


def _encode_image(path: str) -> str:
    with open(path, "rb") as f:
        return base64.standard_b64encode(f.read()).decode("utf-8")


def _image_media_type(path: str) -> str:
    ext = Path(path).suffix.lower()
    mapping = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
               ".webp": "image/webp", ".gif": "image/gif"}
    return mapping.get(ext, "image/png")


class BrandExtractor:
    """
    Extracts brand profile data from various reference materials using
    Claude AI for intelligent interpretation.
    """

    # Source priority hierarchy (lower number = higher priority)
    SOURCE_PRIORITY = {
        "manual": 1,
        "brand_guide_pdf": 2,
        "existing_presentation": 3,
        "website": 4,
        "logo": 5,
        "screenshot": 5,
    }

    def __init__(self, api_key: Optional[str] = None):
        self.client = anthropic.Anthropic(api_key=api_key or os.environ.get("ANTHROPIC_API_KEY"))

    # ------------------------------------------------------------------
    # Public extraction methods
    # ------------------------------------------------------------------

    def from_pdf(self, pdf_path: str, brand_name: str = "Brand") -> BrandProfile:
        """Extract brand profile from a brand guide PDF."""
        try:
            import pdfplumber
        except ImportError:
            raise ImportError("pdfplumber is required for PDF extraction: pip install pdfplumber")

        text_content = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages[:100]:  # cap at 100 pages
                text = page.extract_text() or ""
                text_content.append(text)

        combined_text = "\n\n".join(text_content)
        profile = self._extract_from_text(combined_text, source_type="brand_guide_pdf",
                                          brand_name=brand_name)
        profile.source_materials.append(f"brand_guide_pdf:{Path(pdf_path).name}")
        return profile

    def from_presentation(self, pptx_path: str, brand_name: str = "Brand") -> BrandProfile:
        """Extract brand profile by analyzing an existing .pptx file."""
        try:
            from pptx import Presentation
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
        except ImportError:
            raise ImportError("python-pptx is required: pip install python-pptx")

        prs = Presentation(pptx_path)
        colors_found: list[str] = []
        fonts_found: list[str] = []
        text_samples: list[str] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                text_samples.append(run.text.strip()[:200])
                            if run.font.color and run.font.color.type is not None:
                                try:
                                    rgb = run.font.color.rgb
                                    colors_found.append(_rgb_to_hex(rgb.red, rgb.green, rgb.blue))
                                except Exception:
                                    pass
                            if run.font.name:
                                fonts_found.append(run.font.name)

        # Summarize findings for Claude
        color_counts: dict[str, int] = {}
        for c in colors_found:
            color_counts[c] = color_counts.get(c, 0) + 1
        top_colors = sorted(color_counts, key=lambda x: -color_counts[x])[:10]

        font_counts: dict[str, int] = {}
        for f in fonts_found:
            font_counts[f] = font_counts.get(f, 0) + 1
        top_fonts = sorted(font_counts, key=lambda x: -font_counts[x])[:5]

        summary = (
            f"Brand: {brand_name}\n"
            f"Source: Existing Presentation\n"
            f"Top colors found (hex): {', '.join(top_colors)}\n"
            f"Top fonts found: {', '.join(top_fonts)}\n"
            f"Sample text snippets: {' | '.join(text_samples[:20])}\n"
        )

        profile = self._extract_from_text(summary, source_type="existing_presentation",
                                          brand_name=brand_name)
        profile.source_materials.append(f"existing_presentation:{Path(pptx_path).name}")
        return profile

    def from_image(self, image_path: str, brand_name: str = "Brand",
                   source_type: str = "screenshot") -> BrandProfile:
        """Extract brand profile from a screenshot or logo image."""
        image_data = _encode_image(image_path)
        media_type = _image_media_type(image_path)

        prompt = self._build_image_extraction_prompt(brand_name, source_type)
        profile = self._extract_from_image(image_data, media_type, prompt, brand_name, source_type)
        profile.source_materials.append(f"{source_type}:{Path(image_path).name}")
        return profile

    def from_logo(self, logo_path: str, brand_name: str = "Brand") -> BrandProfile:
        """Extract primary brand colors from a logo file."""
        # Try colorthief first for fast dominant color extraction
        colors_from_thief: list[str] = []
        try:
            from colorthief import ColorThief
            ct = ColorThief(logo_path)
            palette = ct.get_palette(color_count=6, quality=1)
            colors_from_thief = [_rgb_to_hex(*rgb) for rgb in palette]
        except Exception:
            pass

        profile = self.from_image(logo_path, brand_name=brand_name, source_type="logo")

        # Supplement with colorthief results if available
        if colors_from_thief and not profile.colors.primary:
            profile.colors.primary = _make_color_entry(
                colors_from_thief[0], confidence=Confidence.MEDIUM,
                usage="extracted from logo dominant color"
            )
        if colors_from_thief[1:] and not profile.colors.secondary:
            profile.colors.secondary = [
                _make_color_entry(c, confidence=Confidence.LOW) for c in colors_from_thief[1:4]
            ]

        profile.source_materials.append(f"logo:{Path(logo_path).name}")
        return profile

    def from_website(self, url: str, brand_name: str = "Brand") -> BrandProfile:
        """Extract brand profile by scraping a website."""
        try:
            import requests
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("requests and beautifulsoup4 are required: pip install requests beautifulsoup4")

        headers = {"User-Agent": "Mozilla/5.0 (compatible; BrandExtractor/1.0)"}
        resp = requests.get(url, headers=headers, timeout=15)
        resp.raise_for_status()

        soup = BeautifulSoup(resp.text, "html.parser")

        # Extract CSS color values and font families
        style_tags = soup.find_all("style")
        css_text = " ".join(tag.get_text() for tag in style_tags)
        hex_colors = list(set(re.findall(r"#(?:[0-9a-fA-F]{3}){1,2}\b", css_text)))
        font_families = list(set(re.findall(r"font-family\s*:\s*([^;]+)", css_text)))

        # Extract visible text for tone analysis
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        visible_text = soup.get_text(separator=" ", strip=True)[:3000]

        summary = (
            f"Brand: {brand_name}\n"
            f"Website URL: {url}\n"
            f"CSS hex colors found: {', '.join(hex_colors[:20])}\n"
            f"CSS font families: {', '.join(font_families[:5])}\n"
            f"Sample visible text: {visible_text[:2000]}\n"
        )

        profile = self._extract_from_text(summary, source_type="website", brand_name=brand_name)
        profile.source_materials.append(f"website:{urlparse(url).netloc}")
        return profile

    def from_manual(self, data: dict, brand_name: str = "Brand") -> BrandProfile:
        """
        Build a brand profile from manual form input.

        data keys (all optional):
          primary_color, secondary_colors (list), accent_colors (list),
          heading_font, body_font, writing_tone, aspect_ratio
        """
        profile = BrandProfile(name=brand_name)
        profile.source_materials = ["manual"]

        if "primary_color" in data:
            profile.colors.primary = _make_color_entry(
                data["primary_color"], confidence=Confidence.HIGH, usage="manually specified"
            )
        if "secondary_colors" in data:
            profile.colors.secondary = [
                _make_color_entry(c, confidence=Confidence.HIGH) for c in data["secondary_colors"][:4]
            ]
        if "accent_colors" in data:
            profile.colors.accent = [
                _make_color_entry(c, confidence=Confidence.HIGH) for c in data["accent_colors"][:2]
            ]
        if "heading_font" in data:
            profile.typography.heading_font = FontSpec(
                name=data["heading_font"], confidence=Confidence.HIGH
            )
        if "body_font" in data:
            profile.typography.body_font = FontSpec(
                name=data["body_font"], confidence=Confidence.HIGH
            )
        if "writing_tone" in data:
            try:
                profile.content_style.writing_tone = WritingTone(data["writing_tone"])
            except ValueError:
                pass
        if "aspect_ratio" in data:
            try:
                profile.layout.aspect_ratio = AspectRatio(data["aspect_ratio"])
            except ValueError:
                pass

        return profile

    def merge_sources(self, profiles: list[BrandProfile],
                      source_types: list[str]) -> BrandProfile:
        """
        Merge multiple profiles following the priority hierarchy.
        Surfaces conflicts for user review.
        """
        if not profiles:
            return BrandProfile()

        # Sort by priority (ascending = higher priority)
        paired = sorted(
            zip(profiles, source_types),
            key=lambda x: self.SOURCE_PRIORITY.get(x[1], 99),
        )

        merged = paired[0][0]
        for profile, source_type in paired[1:]:
            # Lower priority source only fills in missing fields
            if not merged.colors.primary and profile.colors.primary:
                merged.colors.primary = profile.colors.primary
            if not merged.colors.secondary and profile.colors.secondary:
                merged.colors.secondary = profile.colors.secondary
            if not merged.typography.heading_font and profile.typography.heading_font:
                merged.typography.heading_font = profile.typography.heading_font
            if not merged.typography.body_font and profile.typography.body_font:
                merged.typography.body_font = profile.typography.body_font
            merged.source_materials.extend(profile.source_materials)

        merged.source_materials = list(set(merged.source_materials))
        return merged

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _build_extraction_prompt(self, brand_name: str, source_type: str) -> str:
        return f"""You are a brand design expert. Analyze the provided brand material and extract a structured brand profile.

Brand name: {brand_name}
Source type: {source_type}

Extract the following information and return ONLY a valid JSON object with this exact structure:
{{
  "primary_color": "#RRGGBB or null",
  "secondary_colors": ["#RRGGBB", ...],
  "accent_colors": ["#RRGGBB", ...],
  "background_light": "#RRGGBB or null",
  "background_dark": "#RRGGBB or null",
  "text_heading_color": "#RRGGBB or null",
  "text_body_color": "#RRGGBB or null",
  "heading_font": "font name or null",
  "body_font": "font name or null",
  "size_title": 40,
  "size_h1": 32,
  "size_body": 16,
  "aspect_ratio": "16:9 or 4:3",
  "content_alignment": "left-heavy or centered or grid-based",
  "white_space": "minimal or generous or structured",
  "max_bullets_per_slide": 5,
  "logo_placement": "top-left or top-right or bottom-left or bottom-right",
  "shape_language": "rounded or sharp or geometric or organic",
  "background_treatment": "solid or gradient or textured or image-based",
  "chart_style": "bar or line or donut",
  "writing_tone": "formal or conversational or technical or inspirational",
  "headline_style": "statement or question or metric-led",
  "bullet_style": "full-sentences or fragments",
  "confidences": {{
    "primary_color": "High or Medium or Low",
    "heading_font": "High or Medium or Low",
    "writing_tone": "High or Medium or Low"
  }},
  "review_flags": ["list of fields with Low confidence that need human review"],
  "reasoning": "brief explanation of extraction decisions"
}}

Only include colors you are confident about. Use null for unknown values.
Respond with ONLY the JSON object, no markdown, no explanation."""

    def _build_image_extraction_prompt(self, brand_name: str, source_type: str) -> str:
        return self._build_extraction_prompt(brand_name, source_type)

    def _parse_claude_response(self, raw: str, brand_name: str,
                               source_type: str) -> BrandProfile:
        """Parse Claude's JSON response into a BrandProfile."""
        # Strip markdown fences if present
        raw = raw.strip()
        if raw.startswith("```"):
            raw = re.sub(r"^```[a-z]*\n?", "", raw)
            raw = re.sub(r"\n?```$", "", raw)

        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            # Fallback: return empty profile with review flag
            profile = BrandProfile(name=brand_name)
            profile.review_flags.append("extraction_failed: could not parse AI response")
            return profile

        profile = BrandProfile(name=brand_name)
        confidences = data.get("confidences", {})

        def conf(key: str) -> Confidence:
            val = confidences.get(key, "Medium")
            try:
                return Confidence(val)
            except ValueError:
                return Confidence.MEDIUM

        if data.get("primary_color"):
            profile.colors.primary = _make_color_entry(
                data["primary_color"], confidence=conf("primary_color")
            )
        if data.get("secondary_colors"):
            profile.colors.secondary = [
                _make_color_entry(c, confidence=Confidence.MEDIUM)
                for c in data["secondary_colors"][:4]
                if c
            ]
        if data.get("accent_colors"):
            profile.colors.accent = [
                _make_color_entry(c, confidence=Confidence.MEDIUM)
                for c in data["accent_colors"][:2]
                if c
            ]
        if data.get("background_light"):
            profile.colors.background_light = _make_color_entry(data["background_light"])
        if data.get("background_dark"):
            profile.colors.background_dark = _make_color_entry(data["background_dark"])
        if data.get("text_heading_color"):
            profile.colors.text_heading = _make_color_entry(data["text_heading_color"])
        if data.get("text_body_color"):
            profile.colors.text_body = _make_color_entry(data["text_body_color"])

        if data.get("heading_font"):
            profile.typography.heading_font = FontSpec(
                name=data["heading_font"], confidence=conf("heading_font")
            )
        if data.get("body_font"):
            profile.typography.body_font = FontSpec(name=data["body_font"])
        if data.get("size_title"):
            profile.typography.size_title = int(data["size_title"])
        if data.get("size_h1"):
            profile.typography.size_h1 = int(data["size_h1"])
        if data.get("size_body"):
            profile.typography.size_body = int(data["size_body"])

        try:
            profile.layout.aspect_ratio = AspectRatio(data.get("aspect_ratio", "16:9"))
        except ValueError:
            pass
        profile.layout.content_alignment = data.get("content_alignment", "left-heavy")
        try:
            profile.layout.white_space = WhiteSpacePhilosophy(data.get("white_space", "structured"))
        except ValueError:
            pass
        if data.get("max_bullets_per_slide"):
            profile.layout.max_bullets_per_slide = int(data["max_bullets_per_slide"])

        profile.visuals.logo_placement = data.get("logo_placement", "top-left")
        try:
            profile.visuals.shape_language = ShapeLanguage(data.get("shape_language", "rounded"))
        except ValueError:
            pass
        profile.visuals.background_treatment = data.get("background_treatment", "solid")
        profile.visuals.chart_style = data.get("chart_style", "bar")

        try:
            profile.content_style.writing_tone = WritingTone(
                data.get("writing_tone", "formal")
            )
        except ValueError:
            pass
        try:
            profile.content_style.headline_style = HeadlineStyle(
                data.get("headline_style", "statement")
            )
        except ValueError:
            pass
        profile.content_style.bullet_style = data.get("bullet_style", "fragments")

        profile.review_flags = data.get("review_flags", [])
        return profile

    def _extract_from_text(self, text: str, source_type: str,
                            brand_name: str) -> BrandProfile:
        prompt = self._build_extraction_prompt(brand_name, source_type)
        message = self.client.messages.create(
            model="claude-opus-4-6",
            max_tokens=2048,
            messages=[
                {
                    "role": "user",
                    "content": f"{prompt}\n\n---\n\nBrand material text:\n{text[:12000]}",
                }
            ],
        )
        raw = message.content[0].text
        return self._parse_claude_response(raw, brand_name, source_type)

    def _extract_from_image(self, image_data: str, media_type: str,
                             prompt: str, brand_name: str,
                             source_type: str) -> BrandProfile:
        message = self.client.messages.create(
            model="claude-opus-4-6",
            max_tokens=2048,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": image_data,
                            },
                        },
                        {"type": "text", "text": prompt},
                    ],
                }
            ],
        )
        raw = message.content[0].text
        return self._parse_claude_response(raw, brand_name, source_type)
