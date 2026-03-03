"""
Slide Generation Engine

Generates .pptx presentations from a DeckSpec + BrandProfile.
All slides are fully editable in PowerPoint.

Supported slide types (from PRD Section 5.4):
  - title, section_divider, content_text, content_visual,
    data_chart, comparison, timeline, quote, team, thank_you
"""

import os
import io
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData as PptxChartData

from .brand_profile import BrandProfile, ColorEntry, AspectRatio
from .content_input import DeckSpec, SlideContent, SlideType, ChartData


def _rgb(color_entry: Optional[ColorEntry], fallback: tuple = (0, 0, 0)) -> RGBColor:
    if color_entry and color_entry.hex:
        try:
            hex_val = color_entry.hex.lstrip("#")
            r, g, b = int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
            return RGBColor(r, g, b)
        except Exception:
            pass
    return RGBColor(*fallback)


def _set_font(run, name: Optional[str], size_pt: int, bold: bool = False,
              color: Optional[RGBColor] = None) -> None:
    if name:
        run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _set_bg_color(slide, color: RGBColor) -> None:
    """Set slide background to a solid color."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left: float, top: float, width: float, height: float,
                  text: str, font_name: Optional[str], font_size: int,
                  bold: bool = False, color: RGBColor = RGBColor(0, 0, 0),
                  alignment: PP_ALIGN = PP_ALIGN.LEFT,
                  word_wrap: bool = True) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, font_name, font_size, bold, color)


def _add_colored_rect(slide, left: float, top: float, width: float,
                      height: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border


class SlideGenerator:
    """
    Generates a complete .pptx file from a DeckSpec and BrandProfile.
    """

    # Slide dimensions for 16:9 (inches)
    SLIDE_W = 13.333
    SLIDE_H = 7.5

    # Standard 4:3
    SLIDE_W_43 = 10.0
    SLIDE_H_43 = 7.5

    def __init__(self, brand: BrandProfile):
        self.brand = brand
        self._set_dimensions()

    def _set_dimensions(self) -> None:
        if self.brand.layout.aspect_ratio == AspectRatio.STANDARD:
            self.w = self.SLIDE_W_43
            self.h = self.SLIDE_H_43
        else:
            self.w = self.SLIDE_W
            self.h = self.SLIDE_H

        self.ml = self.brand.layout.margin_left
        self.mr = self.brand.layout.margin_right
        self.mt = self.brand.layout.margin_top
        self.mb = self.brand.layout.margin_bottom
        self.content_w = self.w - self.ml - self.mr
        self.content_h = self.h - self.mt - self.mb

    @property
    def _primary_color(self) -> RGBColor:
        return _rgb(self.brand.colors.primary, (0, 70, 127))

    @property
    def _bg_color(self) -> RGBColor:
        return _rgb(self.brand.colors.background_light, (255, 255, 255))

    @property
    def _text_heading_color(self) -> RGBColor:
        return _rgb(self.brand.colors.text_heading, (0, 0, 0))

    @property
    def _text_body_color(self) -> RGBColor:
        return _rgb(self.brand.colors.text_body, (64, 64, 64))

    @property
    def _accent_color(self) -> RGBColor:
        if self.brand.colors.accent:
            return _rgb(self.brand.colors.accent[0], (255, 165, 0))
        return _rgb(self.brand.colors.secondary[0] if self.brand.colors.secondary else None,
                    (100, 100, 200))

    @property
    def _heading_font(self) -> Optional[str]:
        return self.brand.typography.heading_font.name if self.brand.typography.heading_font else None

    @property
    def _body_font(self) -> Optional[str]:
        return self.brand.typography.body_font.name if self.brand.typography.body_font else None

    # ------------------------------------------------------------------
    # Main generation entry point
    # ------------------------------------------------------------------

    def generate(self, deck: DeckSpec, output_path: str) -> str:
        """Generate the .pptx file and return the output path."""
        prs = Presentation()

        # Set slide dimensions
        if self.brand.layout.aspect_ratio == AspectRatio.STANDARD:
            prs.slide_width = Inches(self.SLIDE_W_43)
            prs.slide_height = Inches(self.SLIDE_H_43)
        else:
            prs.slide_width = Inches(self.SLIDE_W)
            prs.slide_height = Inches(self.SLIDE_H)

        blank_layout = prs.slide_layouts[6]  # completely blank layout

        for slide_content in deck.slides:
            slide = prs.slides.add_slide(blank_layout)
            self._apply_slide(slide, slide_content, deck)

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return output_path

    def regenerate_slide(self, pptx_path: str, slide_index: int,
                         new_content: SlideContent,
                         deck: DeckSpec,
                         output_path: Optional[str] = None) -> str:
        """Regenerate a single slide in an existing presentation."""
        prs = Presentation(pptx_path)
        if slide_index >= len(prs.slides):
            raise IndexError(f"Slide index {slide_index} out of range")

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides[slide_index]

        # Remove all existing shapes
        for shape in list(slide.shapes):
            sp = shape.element
            sp.getparent().remove(sp)

        self._apply_slide(slide, new_content, deck)
        out = output_path or pptx_path
        prs.save(out)
        return out

    # ------------------------------------------------------------------
    # Slide type dispatchers
    # ------------------------------------------------------------------

    def _apply_slide(self, slide, content: SlideContent, deck: DeckSpec) -> None:
        # Set background
        _set_bg_color(slide, self._bg_color)

        dispatch = {
            SlideType.TITLE: self._render_title,
            SlideType.SECTION_DIVIDER: self._render_section_divider,
            SlideType.CONTENT_TEXT: self._render_content_text,
            SlideType.CONTENT_VISUAL: self._render_content_visual,
            SlideType.DATA_CHART: self._render_data_chart,
            SlideType.COMPARISON: self._render_comparison,
            SlideType.TIMELINE: self._render_timeline,
            SlideType.QUOTE: self._render_quote,
            SlideType.TEAM: self._render_team,
            SlideType.THANK_YOU: self._render_thank_you,
        }
        renderer = dispatch.get(content.slide_type, self._render_content_text)
        renderer(slide, content)

        # Add speaker notes
        if content.speaker_notes:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = content.speaker_notes

    # ------------------------------------------------------------------
    # Individual slide renderers
    # ------------------------------------------------------------------

    def _render_title(self, slide, content: SlideContent) -> None:
        """Full-width title slide with primary color accent bar."""
        # Accent bar at bottom
        _add_colored_rect(slide, 0, self.h - 0.5, self.w, 0.5, self._primary_color)

        # Optional decorative left bar
        _add_colored_rect(slide, 0, 0, 0.15, self.h - 0.5, self._primary_color)

        # Title
        _add_text_box(
            slide,
            left=0.4, top=self.h * 0.3,
            width=self.w - 0.5, height=1.5,
            text=content.title,
            font_name=self._heading_font,
            font_size=self.brand.typography.size_title,
            bold=True,
            color=self._text_heading_color,
        )

        # Subtitle
        if content.subtitle:
            _add_text_box(
                slide,
                left=0.4, top=self.h * 0.3 + 1.6,
                width=self.w - 0.5, height=0.8,
                text=content.subtitle,
                font_name=self._body_font,
                font_size=self.brand.typography.size_h3,
                color=self._text_body_color,
            )

    def _render_section_divider(self, slide, content: SlideContent) -> None:
        """Full-color background section divider."""
        _set_bg_color(slide, self._primary_color)

        # Title centered
        _add_text_box(
            slide,
            left=self.ml, top=self.h * 0.35,
            width=self.content_w, height=1.5,
            text=content.title,
            font_name=self._heading_font,
            font_size=self.brand.typography.size_h1,
            bold=True,
            color=RGBColor(255, 255, 255),
            alignment=PP_ALIGN.CENTER,
        )

        if content.subtitle:
            _add_text_box(
                slide,
                left=self.ml, top=self.h * 0.35 + 1.6,
                width=self.content_w, height=0.6,
                text=content.subtitle,
                font_name=self._body_font,
                font_size=self.brand.typography.size_body,
                color=RGBColor(200, 200, 200),
                alignment=PP_ALIGN.CENTER,
            )

    def _render_content_text(self, slide, content: SlideContent) -> None:
        """Standard content slide: headline + bullets or paragraph."""
        # Thin accent line at top
        _add_colored_rect(slide, 0, 0, self.w, 0.08, self._primary_color)

        # Title
        _add_text_box(
            slide,
            left=self.ml, top=self.mt + 0.1,
            width=self.content_w, height=0.8,
            text=content.title,
            font_name=self._heading_font,
            font_size=self.brand.typography.size_h1,
            bold=True,
            color=self._text_heading_color,
        )

        # Divider line below title
        _add_colored_rect(slide, self.ml, self.mt + 0.95, self.content_w, 0.03,
                          self._accent_color)

        body_top = self.mt + 1.1
        body_height = self.h - body_top - self.mb

        if content.body_bullets:
            self._add_bullet_list(slide, content.body_bullets, body_top, body_height)
        elif content.body_text:
            _add_text_box(
                slide,
                left=self.ml, top=body_top,
                width=self.content_w, height=body_height,
                text=content.body_text,
                font_name=self._body_font,
                font_size=self.brand.typography.size_body,
                color=self._text_body_color,
            )

    def _render_content_visual(self, slide, content: SlideContent) -> None:
        """Split layout: text on left, image placeholder on right."""
        _add_colored_rect(slide, 0, 0, self.w, 0.08, self._primary_color)

        half_w = self.content_w / 2 - 0.1

        # Title (full width)
        _add_text_box(
            slide,
            left=self.ml, top=self.mt + 0.1,
            width=self.content_w, height=0.8,
            text=content.title,
            font_name=self._heading_font,
            font_size=self.brand.typography.size_h1,
            bold=True,
            color=self._text_heading_color,
        )

        body_top = self.mt + 1.0

        # Left: bullets / text
        if content.body_bullets:
            self._add_bullet_list(slide, content.body_bullets, body_top,
                                  self.h - body_top - self.mb,
                                  width=half_w)
        elif content.body_text:
            _add_text_box(
                slide,
                left=self.ml, top=body_top,
                width=half_w, height=self.h - body_top - self.mb,
                text=content.body_text,
                font_name=self._body_font,
                font_size=self.brand.typography.size_body,
                color=self._text_body_color,
            )

        # Right: image placeholder box
        img_left = self.ml + half_w + 0.2
        img_top = body_top
        img_w = half_w
        img_h = self.h - body_top - self.mb - 0.2

        _add_colored_rect(slide, img_left, img_top, img_w, img_h,
                          RGBColor(230, 230, 230))

        hint_text = content.image_hint or "[ Image ]"
        _add_text_box(
            slide,
            left=img_left, top=img_top + img_h / 2 - 0.3,
            width=img_w, height=0.6,
            text=hint_text,
            font_name=self._body_font,
            font_size=11,
            color=RGBColor(150, 150, 150),
            alignment=PP_ALIGN.CENTER,
        )

    def _render_data_chart(self, slide, content: SlideContent) -> None:
        """Data/chart slide with title, chart, and takeaway callout."""
        _add_colored_rect(slide, 0, 0, self.w, 0.08, self._primary_color)

        _add_text_box(
            slide,
            left=self.ml, top=self.mt + 0.1,
            width=self.content_w, height=0.7,
            text=content.title,
            font_name=self._heading_font,
            font_size=self.brand.typography.size_h1,
            bold=True,
            color=self._text_heading_color,
        )

        chart_top = self.mt + 0.9
        chart_h = self.h - chart_top - self.mb - 1.0  # leave room for takeaway

        cd = content.chart_data
        if cd and cd.labels and cd.series:
            try:
                self._add_chart(slide, cd, self.ml, chart_top,
                                self.content_w, chart_h)
            except Exception:
                # Chart failed — show placeholder
                _add_colored_rect(slide, self.ml, chart_top, self.content_w, chart_h,
                                  RGBColor(240, 240, 240))
                _add_text_box(slide, self.ml, chart_top + chart_h / 2 - 0.3,
                              self.content_w, 0.6,
                              "[Chart placeholder]", self._body_font, 12,
                              color=RGBColor(150, 150, 150),
                              alignment=PP_ALIGN.CENTER)
        else:
            # Placeholder
            _add_colored_rect(slide, self.ml, chart_top, self.content_w, chart_h,
                              RGBColor(240, 240, 240))
            _add_text_box(slide, self.ml, chart_top + chart_h / 2 - 0.3,
                          self.content_w, 0.6,
                          "[Chart placeholder]", self._body_font, 12,
                          color=RGBColor(150, 150, 150),
                          alignment=PP_ALIGN.CENTER)

        # Takeaway callout at bottom
        takeaway = cd.takeaway if (cd and cd.takeaway) else ""
        if takeaway:
            _add_colored_rect(slide, self.ml, self.h - self.mb - 0.85,
                              self.content_w, 0.7, self._accent_color)
            _add_text_box(
                slide,
                left=self.ml + 0.15,
                top=self.h - self.mb - 0.8,
                width=self.content_w - 0.3,
                height=0.6,
                text=f"Key Takeaway: {takeaway}",
                font_name=self._body_font,
                font_size=self.brand.typography.size_body,
                bold=True,
                color=RGBColor(255, 255, 255),
            )

    def _render_comparison(self, slide, content: SlideContent) -> None:
        """Two-column comparison layout."""
        _add_colored_rect(slide, 0, 0, self.w, 0.08, self._primary_color)

        _add_text_box(
            slide,
            left=self.ml, top=self.mt + 0.1,
            width=self.content_w, height=0.7,
            text=content.title,
            font_name=self._heading_font,
            font_size=self.brand.typography.size_h1,
            bold=True,
            color=self._text_heading_color,
        )

        body_top = self.mt + 1.0
        col_w = (self.content_w - 0.2) / 2
        col_h = self.h - body_top - self.mb

        # Split bullets into two halves
        bullets = content.body_bullets
        mid = len(bullets) // 2
        left_bullets = bullets[:mid] if bullets else []
        right_bullets = bullets[mid:] if bullets else []

        # Left column with colored header
        _add_colored_rect(slide, self.ml, body_top, col_w, 0.4, self._primary_color)
        _add_text_box(slide, self.ml + 0.1, body_top + 0.05, col_w - 0.2, 0.35,
                      "Option A", self._heading_font, 14, bold=True,
                      color=RGBColor(255, 255, 255))
        self._add_bullet_list(slide, left_bullets, body_top + 0.45,
                              col_h - 0.5, left=self.ml, width=col_w)

        # Right column with colored header
        right_left = self.ml + col_w + 0.2
        _add_colored_rect(slide, right_left, body_top, col_w, 0.4, self._accent_color)
        _add_text_box(slide, right_left + 0.1, body_top + 0.05, col_w - 0.2, 0.35,
                      "Option B", self._heading_font, 14, bold=True,
                      color=RGBColor(255, 255, 255))
        self._add_bullet_list(slide, right_bullets, body_top + 0.45,
                              col_h - 0.5, left=right_left, width=col_w)

    def _render_timeline(self, slide, content: SlideContent) -> None:
        """Horizontal timeline slide."""
        _add_colored_rect(slide, 0, 0, self.w, 0.08, self._primary_color)

        _add_text_box(
            slide,
            left=self.ml, top=self.mt + 0.1,
            width=self.content_w, height=0.7,
            text=content.title,
            font_name=self._heading_font,
            font_size=self.brand.typography.size_h1,
            bold=True,
            color=self._text_heading_color,
        )

        bullets = content.body_bullets or ["Phase 1", "Phase 2", "Phase 3"]
        n = len(bullets)
        timeline_y = self.h / 2

        # Horizontal line
        _add_colored_rect(slide, self.ml, timeline_y - 0.03,
                          self.content_w, 0.06, self._primary_color)

        # Nodes and labels
        step_w = self.content_w / n
        for i, label in enumerate(bullets):
            cx = self.ml + step_w * i + step_w / 2

            # Circle marker
            marker_size = 0.25
            _add_colored_rect(slide, cx - marker_size / 2,
                              timeline_y - marker_size / 2,
                              marker_size, marker_size,
                              self._accent_color if i % 2 == 0 else self._primary_color)

            # Label above or below alternating
            label_y = timeline_y - 1.0 if i % 2 == 0 else timeline_y + 0.35
            _add_text_box(
                slide,
                left=cx - step_w / 2,
                top=label_y,
                width=step_w,
                height=0.6,
                text=label,
                font_name=self._body_font,
                font_size=self.brand.typography.size_caption,
                color=self._text_body_color,
                alignment=PP_ALIGN.CENTER,
            )

    def _render_quote(self, slide, content: SlideContent) -> None:
        """Large centered quote slide."""
        _set_bg_color(slide, _rgb(self.brand.colors.background_dark, (30, 30, 30))
                      if self.brand.colors.background_dark else RGBColor(30, 30, 30))

        quote_text = content.body_text or content.title
        _add_text_box(
            slide,
            left=self.ml, top=self.h * 0.2,
            width=self.content_w, height=self.h * 0.5,
            text=f"\u201c{quote_text}\u201d",
            font_name=self._heading_font,
            font_size=28,
            bold=True,
            color=RGBColor(255, 255, 255),
            alignment=PP_ALIGN.CENTER,
        )

        if content.subtitle:
            _add_text_box(
                slide,
                left=self.ml, top=self.h * 0.72,
                width=self.content_w, height=0.5,
                text=f"\u2014 {content.subtitle}",
                font_name=self._body_font,
                font_size=self.brand.typography.size_body,
                color=self._accent_color,
                alignment=PP_ALIGN.CENTER,
            )

    def _render_team(self, slide, content: SlideContent) -> None:
        """Team grid slide."""
        _add_colored_rect(slide, 0, 0, self.w, 0.08, self._primary_color)

        _add_text_box(
            slide,
            left=self.ml, top=self.mt + 0.1,
            width=self.content_w, height=0.7,
            text=content.title,
            font_name=self._heading_font,
            font_size=self.brand.typography.size_h1,
            bold=True,
            color=self._text_heading_color,
        )

        members = content.body_bullets or ["Team Member"]
        n = min(len(members), 6)
        cols = min(n, 3)
        rows = (n + cols - 1) // cols
        card_w = self.content_w / cols - 0.15
        card_h = (self.h - self.mt - 1.1 - self.mb) / rows - 0.1

        for i, member in enumerate(members[:6]):
            row = i // cols
            col = i % cols
            card_left = self.ml + col * (card_w + 0.15)
            card_top = self.mt + 1.1 + row * (card_h + 0.1)

            # Avatar placeholder
            avatar_size = min(card_w * 0.4, card_h * 0.55)
            _add_colored_rect(slide, card_left + (card_w - avatar_size) / 2,
                              card_top, avatar_size, avatar_size,
                              self._primary_color)

            # Name / bio text
            _add_text_box(
                slide,
                left=card_left, top=card_top + avatar_size + 0.1,
                width=card_w, height=card_h - avatar_size - 0.15,
                text=member,
                font_name=self._body_font,
                font_size=11,
                color=self._text_body_color,
                alignment=PP_ALIGN.CENTER,
            )

    def _render_thank_you(self, slide, content: SlideContent) -> None:
        """Closing / CTA slide."""
        _add_colored_rect(slide, 0, 0, self.w, 0.08, self._primary_color)
        _add_colored_rect(slide, 0, self.h - 0.08, self.w, 0.08, self._primary_color)

        main_text = content.title or "Thank You"
        _add_text_box(
            slide,
            left=self.ml, top=self.h * 0.28,
            width=self.content_w, height=1.2,
            text=main_text,
            font_name=self._heading_font,
            font_size=self.brand.typography.size_title,
            bold=True,
            color=self._text_heading_color,
            alignment=PP_ALIGN.CENTER,
        )

        if content.subtitle:
            _add_text_box(
                slide,
                left=self.ml, top=self.h * 0.28 + 1.3,
                width=self.content_w, height=0.6,
                text=content.subtitle,
                font_name=self._body_font,
                font_size=self.brand.typography.size_body,
                color=self._text_body_color,
                alignment=PP_ALIGN.CENTER,
            )

        if content.body_bullets:
            contact_text = " | ".join(content.body_bullets)
            _add_text_box(
                slide,
                left=self.ml, top=self.h * 0.65,
                width=self.content_w, height=0.5,
                text=contact_text,
                font_name=self._body_font,
                font_size=self.brand.typography.size_caption,
                color=self._text_body_color,
                alignment=PP_ALIGN.CENTER,
            )

    # ------------------------------------------------------------------
    # Shared rendering helpers
    # ------------------------------------------------------------------

    def _add_bullet_list(self, slide, bullets: list[str], top: float,
                         height: float, left: Optional[float] = None,
                         width: Optional[float] = None) -> None:
        left = left if left is not None else self.ml
        width = width if width is not None else self.content_w

        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        txBox.word_wrap = True
        tf = txBox.text_frame
        tf.word_wrap = True

        max_bullets = self.brand.layout.max_bullets_per_slide
        bullets = bullets[:max_bullets]

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.level = 0

            # Bullet character
            run = p.add_run()
            run.text = f"\u2022  {bullet}"
            _set_font(run, self._body_font,
                      self.brand.typography.size_body,
                      color=self._text_body_color)

            # Paragraph spacing
            from pptx.util import Pt as Pt_
            p.space_before = Pt_(4)
            p.space_after = Pt_(4)

    def _add_chart(self, slide, cd: ChartData,
                   left: float, top: float,
                   width: float, height: float) -> None:
        """Add a native PowerPoint chart to the slide."""
        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "donut": XL_CHART_TYPE.DOUGHNUT,
            "pie": XL_CHART_TYPE.PIE,
        }
        xl_type = chart_type_map.get(cd.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_data = PptxChartData()
        chart_data.categories = cd.labels

        for series in cd.series:
            chart_data.add_series(series.get("name", "Series"),
                                  series.get("values", []))

        chart_shape = slide.shapes.add_chart(
            xl_type,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data,
        )

        chart = chart_shape.chart
        chart.has_legend = len(cd.series) > 1
        if chart.has_title:
            chart.chart_title.text_frame.text = cd.title
